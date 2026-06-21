#!/usr/bin/env python3
"""
selftest_stage2.py - verifies Stage 2 (PDF/DOCX/XLSX text extraction).

For each format, if a library that can WRITE a sample is installed, it generates a
tiny file with known content and checks the extractor produces the expected units
(with the extracted `text` carried along). Formats whose libs are absent are skipped
with a note. Independently, it always checks GRACEFUL DEGRADATION: with the reader
library hidden, the chunker returns None (-> the file is skipped, never a crash).

Run:  python selftest_stage2.py
"""
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))
import extract_structure as es


def _have(mod: str) -> bool:
    try:
        __import__(mod); return True
    except Exception:
        return False


def _fallback_returns_none(reader_mod: str, fn, path: Path) -> bool:
    """Hide the reader library -> the chunker must return None, not raise."""
    saved = sys.modules.get(reader_mod, "absent")
    sys.modules[reader_mod] = None                     # makes `import reader_mod` fail
    try:
        return fn(path, "x." + path.suffix.lstrip("."), "absorb") is None
    finally:
        if saved == "absent":
            sys.modules.pop(reader_mod, None)
        else:
            sys.modules[reader_mod] = saved


def main() -> int:
    tmp = Path(tempfile.mkdtemp(prefix="amg-s2-"))
    tested = []

    # classification routing (independent of libraries)
    assert es.classify(Path("a.pdf"))[:2] == ("doc", "pdf")
    assert es.classify(Path("a.docx"))[:2] == ("doc", "docx")
    assert es.classify(Path("a.xlsx"))[:2] == ("data", "xlsx")
    assert es.classify(Path("a.pptx"))[:2] == ("doc", "pptx")   # now extracted, not skipped
    assert es.classify(Path("a.doc"))[1] == "skip"      # legacy format stays skipped
    assert es.classify(Path("a.ppt"))[1] == "skip"      # legacy .ppt stays skipped
    print("PASS  classify: pdf->doc, docx->doc, xlsx->data, pptx->doc, .doc/.ppt->skip")

    # ---- PDF (write with fpdf2, read with pypdf) ----------------------------
    if _have("fpdf") and _have("pypdf"):
        from fpdf import FPDF
        pdf = FPDF(); pdf.add_page(); pdf.set_font("helvetica", size=12)
        pdf.multi_cell(0, 10, "Routing conventions and request maps.")
        pdf.add_page(); pdf.multi_cell(0, 10, "Second page about controllers.")
        p = tmp / "doc.pdf"; pdf.output(str(p))
        units = es._pdf_units(p, "doc.pdf", "absorb")
        assert units and all(u["category"] == "doc" and u["text"] for u in units), units
        assert units[0]["id"] == "doc:doc.pdf::p1" and units[0]["kind"] == "page"
        assert _fallback_returns_none("pypdf", es._pdf_units, p)
        print(f"PASS  pdf: {len(units)} page unit(s) with text; missing pypdf -> skip")
        tested.append("pdf")
    else:
        print("SKIP  pdf (need: pip install fpdf2 pypdf)")

    # ---- DOCX (python-docx reads and writes) --------------------------------
    if _have("docx"):
        import docx
        d = docx.Document()
        d.add_heading("Маршрутизация", level=1)
        d.add_paragraph("Как запросы попадают в контроллеры.")
        d.add_heading("Конфигурация", level=1)
        d.add_paragraph("Переменные окружения.")
        p = tmp / "guide.docx"; d.save(str(p))
        units = es._docx_units(p, "guide.docx", "absorb")
        quals = {u["qualname"] for u in units}
        assert units and all(u["category"] == "doc" and u["text"] for u in units), units
        assert any("section" == u["kind"] for u in units) and len(quals) >= 2, quals
        assert _fallback_returns_none("docx", es._docx_units, p)
        print(f"PASS  docx: {len(units)} heading-section unit(s) with text; missing lib -> skip")
        tested.append("docx")
    else:
        print("SKIP  docx (need: pip install python-docx)")

    # ---- XLSX (openpyxl reads and writes) -----------------------------------
    if _have("openpyxl"):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active; ws.title = "Routes"
        ws.append(["path", "handler", "method"])
        ws.append(["/users", "list_users", "GET"])
        ws.append(["/users", "create_user", "POST"])
        wb.create_sheet("Config").append(["key", "value"])
        p = tmp / "data.xlsx"; wb.save(str(p))
        units = es._xlsx_units(p, "data.xlsx", "absorb")
        assert len(units) == 2, units                   # one unit per sheet
        u0 = units[0]
        assert u0["category"] == "data" and u0["kind"] == "sheet"
        assert u0["id"] == "data:data.xlsx::Routes" and "Columns:" in u0["text"]
        assert _fallback_returns_none("openpyxl", es._xlsx_units, p)
        print(f"PASS  xlsx: {len(units)} sheet unit(s), structural text; missing lib -> skip")
        tested.append("xlsx")
    else:
        print("SKIP  xlsx (need: pip install openpyxl)")

    # ---- PPTX (python-pptx reads and writes; Stage 11) ----------------------
    if _have("pptx"):
        from pptx import Presentation
        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[5])     # "Title Only" layout
        s1.shapes.title.text = "Routing overview"
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "Controllers"
        p = tmp / "deck.pptx"; prs.save(str(p))
        units = es._pptx_units(p, "deck.pptx", "absorb")
        assert units and all(u["category"] == "doc" and u["text"] for u in units), units
        assert units[0]["id"] == "doc:deck.pptx::s1" and units[0]["kind"] == "section", units[0]
        assert _fallback_returns_none("pptx", es._pptx_units, p)
        print(f"PASS  pptx: {len(units)} slide unit(s) with text; missing python-pptx -> skip")
        tested.append("pptx")
    else:
        print("SKIP  pptx (need: pip install python-pptx)")

    # ---- tree-sitter (canonical kinds for non-Python code; stage 1, task 8) --
    if _have("tree_sitter_language_pack"):
        p = tmp / "app.js"
        p.write_text("function foo(a) {\n  return bar(a);\n}\n\n"
                     "class Baz {\n  run() {\n    return foo(1);\n  }\n}\n",
                     encoding="utf-8")
        units = es._treesitter_units(p, "app.js", "absorb", "javascript")
        assert units is not None, "javascript grammar failed to load"
        kinds = {u["qualname"]: u["kind"] for u in units if u["qualname"]}
        assert kinds.get("foo") == "function" and kinds.get("Baz") == "class", kinds
        assert all(u["kind"] in ("module", "function", "class") for u in units), units
        # Stage 13: line_end is the unit's real end line (multi-line span), not lineno
        spans = {u["qualname"]: (u["lineno"], u["line_end"]) for u in units if u["qualname"]}
        assert spans["foo"] == (1, 3), spans                # function foo spans lines 1-3
        assert spans["Baz"][1] > spans["Baz"][0], spans     # class Baz spans several lines
        print("PASS  treesitter: grammar kinds canonicalized to function/class; line_end spans")
        tested.append("treesitter")
    else:
        print("SKIP  treesitter (need: pip install tree-sitter tree-sitter-language-pack)")

    print(f"\nSTAGE 2 CHECKS PASSED (formats exercised: {', '.join(tested) or 'none — install libs'})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
