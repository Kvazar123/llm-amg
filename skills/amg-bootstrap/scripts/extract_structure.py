#!/usr/bin/env python3
"""
extract_structure.py - deterministic source -> structural units. (Stage 1: unified input)

The graph engine is domain-blind: it stores nodes/edges/text and does not care whether
a unit came from code, prose, or data. Only THIS file is type-aware. It is the "sensory
cortex" of the system: it routes each file to the right chunker by modality, then the
single associative graph is built from the result.

Pipeline per file:
  1. ignore?      built-in defaults (caches, deps, binaries) + the repo's .gitignore
  2. classify     by extension, then a content sniff for the ambiguous/extensionless
  3. chunk        a registry maps (chunker -> how to split into content-hashed units):
       python      stdlib `ast`  -> module + functions/classes/methods, + imports + CALLS
       <lang>code  tree-sitter   -> functions/classes + calls   (OPTIONAL; see below)
                   if tree-sitter or the grammar is unavailable -> one unit per file
       markdown    headings      -> one unit per section
       rst         adornments    -> one unit per underline-headed section
       text        paragraphs    -> one unit per blank-line block
       log         episodes      -> one block per bounded window of timestamped lines
       json/yaml   records       -> one unit per entry; large nests split by key path
       ndjson      records       -> one unit per JSON line
       (json/ndjson that look like a chat export -> one section per message + follows)
       csv/tsv     table         -> one structural unit (headers + sample rows)
       pdf/docx/xlsx/pptx         -> page / section / sheet / slide   (OPTIONAL libs)
  4. tag          each unit carries `category` (code|doc|data) -> physical node bucket,
                  and `policy` (mirror|absorb) inherited from its source.

tree-sitter is OPTIONAL. Python is handled fully by the stdlib `ast`, with NO dependency.
Other languages get function-level granularity + call edges ONLY if
`tree-sitter-language-pack` is importable and the grammar loads; otherwise the file
becomes a single unit (coarser, but never an error). Install it to get the most:
    pip install tree-sitter tree-sitter-language-pack

The content hash is what makes the whole system idempotent: reconcile.py compares these
hashes to the graph to decide what changed.

CLI:
    python extract_structure.py <project_root> [<amg_root>]            # dump units JSON
    python extract_structure.py <project_root> [<amg_root>] --stats     # human summary
"""
from __future__ import annotations

import ast
import fnmatch
import hashlib
import json
import re
import sys
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Optional, Set, Tuple

try:
    import yaml
except ImportError:                       # pragma: no cover
    sys.stderr.write("extract_structure.py needs PyYAML: pip install pyyaml\n")
    raise

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[union-attr]
except (AttributeError, ValueError):
    pass

# --------------------------------------------------------------------------- #
# Classification tables
# --------------------------------------------------------------------------- #

# Directories never indexed (hygiene, not meaning): a cheap pre-semantic filter,
# like sensory gating. Both common agent dirs (.claude / .agents) are here so the
# memory never indexes itself; the CONFIGURED agent dir is added on top at run time
# (_effective_ignore_dirs), covering a custom name. The repo's .gitignore stacks too.
DEFAULT_IGNORE_DIRS = {
    ".git", ".hg", ".svn", ".claude", ".agents", "node_modules", ".venv", "venv", "env",
    "__pycache__", ".pytest_cache", ".mypy_cache", ".ruff_cache", ".tox",
    "dist", "build", "target", "out", ".next", ".nuxt", ".cache", "vendor",
    "site-packages", ".idea", ".vscode", ".gradle", "coverage", ".terraform",
    "bin", "obj",
}

BINARY_EXT = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".ico", ".webp", ".tif", ".tiff",
    ".doc", ".xls", ".ppt",                 # legacy office formats: no pure-Python extractor
    ".zip", ".tar", ".gz", ".bz2", ".7z", ".rar", ".jar", ".war",
    ".exe", ".dll", ".so", ".dylib", ".bin", ".pyc", ".pyo", ".class", ".o", ".a",
    ".woff", ".woff2", ".ttf", ".eot", ".otf",
    ".mp3", ".mp4", ".avi", ".mov", ".wav", ".flac", ".ogg", ".webm",
    ".db", ".sqlite", ".sqlite3", ".pdb", ".lock", ".svg",
}

# extension -> tree-sitter grammar name (for non-Python code)
CODE_LANG_BY_EXT = {
    ".js": "javascript", ".jsx": "javascript", ".mjs": "javascript", ".cjs": "javascript",
    ".ts": "typescript", ".tsx": "tsx",
    ".php": "php", ".c": "c", ".h": "c", ".cc": "cpp", ".cpp": "cpp", ".cxx": "cpp",
    ".hpp": "cpp", ".hh": "cpp", ".java": "java", ".go": "go", ".rs": "rust",
    ".rb": "ruby", ".cs": "c_sharp", ".swift": "swift", ".kt": "kotlin",
    ".scala": "scala", ".lua": "lua", ".sh": "bash", ".bash": "bash", ".sql": "sql",
    ".pl": "perl", ".r": "r", ".dart": "dart", ".ex": "elixir", ".exs": "elixir",
}

PROSE_EXT = {".md": "headings", ".mdx": "headings", ".markdown": "headings",
             ".txt": "paragraphs", ".text": "paragraphs"}

DATA_EXT = {".json": "json", ".yaml": "json", ".yml": "json"}


def _sha(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


# --------------------------------------------------------------------------- #
# Config + source resolution (mirror_path / absorb_path, str or list)
# --------------------------------------------------------------------------- #

def load_config(amg_root: Path) -> Dict[str, Any]:
    return yaml.safe_load((amg_root / "config.yml").read_text(encoding="utf-8")) or {}


def _as_list(val: Any) -> List[str]:
    if val is None:
        return []
    if isinstance(val, str):
        return [val]
    return [str(v) for v in val]


def resolve_sources(config: Dict[str, Any]) -> List[Tuple[str, str]]:
    """Return [(path, policy)]. Prefers mirror_path/absorb_path/absorb_once_path; falls
    back to a legacy `sources:` dict so existing configs keep working. Order matters: on
    an overlap the dedup in reconcile.plan keeps the LAST, so the most-preserving policy
    (absorb_once, then absorb, then mirror) wins."""
    out: List[Tuple[str, str]] = []
    if any(k in config for k in ("mirror_path", "absorb_path", "absorb_once_path")):
        out += [(p, "mirror") for p in _as_list(config.get("mirror_path"))]
        out += [(p, "absorb") for p in _as_list(config.get("absorb_path"))]
        out += [(p, "absorb_once") for p in _as_list(config.get("absorb_once_path"))]
        return out
    # legacy form: sources: {name: {path, policy}}
    for _, src in (config.get("sources") or {}).items():
        if isinstance(src, dict) and src.get("path"):
            out.append((src["path"], src.get("policy", "mirror")))
    return out


def detect_policy_conflicts(units: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Unit ids produced under MORE THAN ONE policy — a file that falls under two source
    roots of different intent (e.g. mirror_path: . and absorb_path: data; audit 1.29).
    The dedup in reconcile.plan keeps the LAST source in resolve_sources order
    (mirror -> absorb -> absorb_once), so the most-preserving policy wins; this surfaces
    the overlap so the choice is never silent. Returns [{id, policies:[...]}], sorted."""
    pols: Dict[str, Set[str]] = {}
    for u in units:
        p = u.get("policy")
        if p:
            pols.setdefault(u["id"], set()).add(p)
    return [{"id": uid, "policies": sorted(ps)}
            for uid, ps in sorted(pols.items()) if len(ps) > 1]


# --------------------------------------------------------------------------- #
# Ignore (.gitignore + defaults) and file walking
# --------------------------------------------------------------------------- #

def load_gitignore(project_root: Path) -> List[str]:
    f = project_root / ".gitignore"
    if not f.exists():
        return []
    pats = []
    for line in f.read_text(encoding="utf-8", errors="replace").splitlines():
        line = line.strip()
        if line and not line.startswith("#") and not line.startswith("!"):
            pats.append(line.rstrip("/"))
    return pats


def _gitignored(rel: str, patterns: List[str]) -> bool:
    base = rel.split("/")[-1]
    segs = rel.split("/")
    for pat in patterns:
        if fnmatch.fnmatch(rel, pat) or fnmatch.fnmatch(base, pat):
            return True
        if pat in segs:                       # a directory name anywhere in the path
            return True
        if fnmatch.fnmatch(rel, pat + "/*"):
            return True
    return False


def _norm_globs(vals: Any) -> List[str]:
    """Normalize an exclude list the way gitignore patterns are read: drop blanks and a
    trailing '/' so 'raw/scratch/' matches like the gitignore rule would (via the
    pat + '/*' branch in _gitignored). Used for `exclude` and the per-intent variants."""
    return [v.rstrip("/") for v in _as_list(vals) if str(v).strip()]


def _excludes_for_policy(config: Dict[str, Any], policy: str) -> List[str]:
    """Global `exclude` plus the per-intent `mirror_exclude` / `absorb_exclude` for this
    source's policy (additive). All are glob patterns matched by the same engine as
    .gitignore. Empty by default, so behavior is unchanged until the keys are set."""
    out = _norm_globs(config.get("exclude"))
    key = {"mirror": "mirror_exclude", "absorb": "absorb_exclude",
           "absorb_once": "absorb_exclude"}.get(policy)
    if key:
        out = out + _norm_globs(config.get(key))
    return out


def _effective_ignore_dirs(amg_root: Optional[Path]) -> Set[str]:
    """DEFAULT_IGNORE_DIRS plus the RESOLVED agent dir, so the engine never indexes its
    own directory whatever it is named (.claude / .agents / custom; roadmap 4.9). The
    name is derived from the store location <agent_dir>/amg, which is authoritative —
    independent of the (advisory) `agent_dir` config key."""
    dirs = set(DEFAULT_IGNORE_DIRS)
    if amg_root is not None:
        dirs.add(Path(amg_root).resolve().parent.name)
    return dirs


def _gitignore_for_source(base_rel: str, gitignore: List[str]) -> List[str]:
    """The .gitignore patterns that still apply under an EXPLICITLY configured source
    `base_rel`. A pattern matching the source ROOT itself is dropped: naming a path in
    mirror_path/absorb_path is a deliberate opt-in, so it is ingested even when
    .gitignore lists it (e.g. absorb_path: logs where .gitignore has logs/). Patterns
    matching only deeper junk still apply, so a wide source (mirror_path: .) keeps full
    .gitignore hygiene. DEFAULT_IGNORE_DIRS is a separate hard layer and is NOT relaxed
    here — the agent dir must never be indexed."""
    norm = base_rel.strip("/").replace("\\", "/")
    if not norm or norm == ".":
        return gitignore                      # whole-project source: keep every pattern
    return [p for p in gitignore if not _gitignored(norm, [p])]


def iter_source_files(project_root: Path, base_rel: str, extra_excludes: List[str],
                      gitignore: List[str], ignore_dirs: Set[str]) -> Iterable[Path]:
    """Yield ingestible files under a configured source. Filtering, in order: the hard
    DEFAULT_IGNORE / agent-dir set (ignore_dirs), then .gitignore minus the patterns the
    explicit source opted past, then the source's effective excludes."""
    base = (project_root / base_rel).resolve()
    if not base.exists():
        return
    src_gitignore = _gitignore_for_source(base_rel, gitignore)
    for p in base.rglob("*"):
        if not p.is_file():
            continue
        if any(part in ignore_dirs for part in p.parts):
            continue
        rel = p.relative_to(project_root).as_posix()
        if _gitignored(rel, src_gitignore) or _gitignored(rel, extra_excludes):
            continue
        yield p


# --------------------------------------------------------------------------- #
# Classifier: (category, chunker, lang, ambiguous)
# --------------------------------------------------------------------------- #

def classify(path: Path) -> Tuple[str, str, Optional[str], bool]:
    ext = path.suffix.lower()
    if ext in BINARY_EXT:
        return ("binary", "skip", None, False)
    if ext == ".py":
        return ("code", "python", "python", False)
    if ext in CODE_LANG_BY_EXT:
        return ("code", "treesitter", CODE_LANG_BY_EXT[ext], False)
    if ext in PROSE_EXT:
        return ("doc", PROSE_EXT[ext], None, False)
    if ext == ".rst":
        return ("doc", "rst", None, False)            # reStructuredText underline headings
    if ext in DATA_EXT:
        return ("data", "json", None, False)
    if ext == ".ndjson":
        return ("data", "ndjson", None, False)         # newline-delimited JSON (one obj/line)
    if ext in (".csv", ".tsv"):
        return ("data", "csv", None, False)
    if ext == ".log":
        return ("doc", "log", None, False)             # timestamped log -> episode blocks
    if ext == ".pdf":
        return ("doc", "pdf", "pdf", False)
    if ext == ".docx":
        return ("doc", "docx", "docx", False)
    if ext == ".pptx":
        return ("doc", "pptx", "pptx", False)
    if ext in (".xlsx", ".xlsm"):
        return ("data", "xlsx", "xlsx", False)
    # extensionless or unknown extension -> sniff content
    try:
        head = path.read_bytes()[:2048]
    except OSError:
        return ("binary", "skip", None, False)
    if b"\x00" in head:
        return ("binary", "skip", None, False)
    # readable text of unknown kind: treat as prose, flag as ambiguous so the
    # bootstrap skill may ask the classifier subagent to refine it.
    return ("doc", "paragraphs", None, True)


# --------------------------------------------------------------------------- #
# Classifier overrides: amg-classifier's verdict made effective in code (1.13)
# --------------------------------------------------------------------------- #

# Categories an override may assign (the amg-classifier subagent picks one of these).
_OVERRIDE_CATEGORIES = {"code", "doc", "data"}


def load_overrides(amg_root: Optional[Path]) -> Dict[str, Dict[str, Any]]:
    """Read the amg-classifier category overrides for ambiguous files (audit 1.13).

    Format: { "<rel/path>": {"category": code|doc|data, "language": <grammar|null>} }.
    The bootstrap workflow writes this to work/classification-overrides.json after the
    amg-classifier subagent labels the files extract flagged ambiguous. A missing file
    -> {} (pure deterministic classification); a malformed file warns and is ignored, so
    a bad override never blocks bootstrap.
    """
    if amg_root is None:
        return {}
    f = amg_root / "work" / "classification-overrides.json"
    if not f.exists():
        return {}
    try:
        data = json.loads(f.read_text(encoding="utf-8"))
    except (OSError, ValueError):
        sys.stderr.write(f"warning: ignoring malformed {f}\n")
        return {}
    return data if isinstance(data, dict) else {}


def _route_override(category: str, language: Optional[str]) -> Tuple[str, str, Optional[str]]:
    """Map an amg-classifier override (category + optional grammar name) onto the
    (category, chunker, lang) the chunker registry expects: code+python -> the stdlib
    `ast` chunker; code+grammar -> tree-sitter (a null grammar degrades to a single file
    unit, exactly like any unavailable grammar); data -> the json/yaml record chunker;
    doc -> the prose paragraph chunker (the same safe default the content sniff uses — a
    'use headings' hint is not expressible, so structured markdown should keep its .md)."""
    if category == "code":
        if language == "python":
            return ("code", "python", "python")
        return ("code", "treesitter", language)
    if category == "data":
        return ("data", "json", None)
    return ("doc", "paragraphs", None)


def _classify_path(path: Path, rel: str, overrides: Dict[str, Dict[str, Any]]
                   ) -> Tuple[str, str, Optional[str], bool, bool]:
    """classify(), but an explicit override for this rel-path wins over the
    deterministic guess (applied BEFORE the content sniff, so it routes the file
    straight to the chosen chunker). Returns classify()'s 4-tuple plus a final
    `overridden` flag; an override resolves ambiguity, so its `ambiguous` is False."""
    ov = overrides.get(rel)
    if isinstance(ov, dict) and ov.get("category") in _OVERRIDE_CATEGORIES:
        category, chunker, lang = _route_override(ov["category"], ov.get("language"))
        return (category, chunker, lang, False, True)
    category, chunker, lang, amb = classify(path)
    return (category, chunker, lang, amb, False)


# --------------------------------------------------------------------------- #
# Chunkers
# --------------------------------------------------------------------------- #

def _file_unit(rel: str, category: str, policy: str, text: str, lang: Optional[str] = None) -> Dict[str, Any]:
    return {"id": f"{category}:{rel}", "kind": "file", "source_path": rel,
            "category": category, "policy": policy, "qualname": "", "lineno": 1,
            "lang": lang, "content_sha": _sha(text)}


def _python_units(path: Path, rel: str, policy: str) -> List[Dict[str, Any]]:
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines(keepends=True)
    try:
        tree = ast.parse(text)
    except SyntaxError:
        return [_file_unit(rel, "code", policy, text, "python")]

    imports: List[str] = []
    for n in ast.walk(tree):
        if isinstance(n, ast.Import):
            imports += [a.name for a in n.names]
        elif isinstance(n, ast.ImportFrom) and n.module:
            imports.append(n.module)

    units = [{"id": f"code:{rel}", "kind": "module", "source_path": rel,
              "category": "code", "policy": policy, "qualname": "", "lineno": 1,
              "lang": "python", "content_sha": _sha(text),
              "imports": sorted(set(imports))}]

    def calls_in(node: ast.AST) -> List[str]:
        names = []
        for sub in ast.walk(node):
            if isinstance(sub, ast.Call):
                f = sub.func
                if isinstance(f, ast.Name):
                    names.append(f.id)
                elif isinstance(f, ast.Attribute):
                    names.append(f.attr)
        return sorted(set(names))

    def slice_src(node: ast.stmt) -> str:
        return "".join(lines[node.lineno - 1: getattr(node, "end_lineno", node.lineno)])

    def walk(node: ast.Module | ast.ClassDef, prefix: str) -> None:
        for child in node.body:
            if isinstance(child, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                qual = f"{prefix}{child.name}"
                units.append({
                    "id": f"code:{rel}::{qual}",
                    "kind": "class" if isinstance(child, ast.ClassDef) else "function",
                    "source_path": rel, "category": "code", "policy": policy,
                    "qualname": qual, "lineno": child.lineno, "lang": "python",
                    "content_sha": _sha(slice_src(child)), "calls": calls_in(child)})
                if isinstance(child, ast.ClassDef):
                    walk(child, prefix=f"{qual}.")

    walk(tree, prefix="")
    return units


# tree-sitter node-type families (vary by grammar; matched broadly), mapped to
# the canonical node types so non-Python code gets the same retrieval tiers and
# path:line pointers as Python (roadmap 1.25; stage 1, task 8). Type-level
# containers (struct/impl/trait/interface/enum) canonicalize to `class`.
_TS_DEF = {
    "function_definition": "function", "function_declaration": "function",
    "method_definition": "function", "method_declaration": "function",
    "function_item": "function", "constructor_declaration": "function",
    "class_definition": "class", "class_declaration": "class",
    "class_specifier": "class", "struct_specifier": "class", "struct_item": "class",
    "impl_item": "class", "trait_item": "class",
    "interface_declaration": "class", "enum_declaration": "class",
}
_TS_CALL = {"call", "call_expression", "function_call_expression",
            "method_invocation", "invocation_expression"}


def _treesitter_units(path: Path, rel: str, policy: str, lang: str) -> Optional[List[Dict[str, Any]]]:
    """Function/class granularity + calls via tree-sitter. Returns None to signal
    'unavailable -> fall back to a single file unit'.

    Two binding generations ship under the same package name: the classic pack
    mirrors py-tree-sitter (parse(bytes); node.type/.children/.text properties;
    start_point), the alef rewrite (>= 1.8) exposes a method-based API
    (parse(str); node.kind()/child(i); byte offsets; start_position()). The
    grammar node-type strings are identical, so only access is feature-detected.
    """
    try:
        from tree_sitter_language_pack import get_parser
        parser = get_parser(lang)
    except Exception:
        return None
    try:
        data = path.read_bytes()
        try:
            tree = parser.parse(data)         # type: ignore[arg-type]
        except TypeError:                     # alef binding wants str, not bytes
            tree = parser.parse(data.decode("utf-8", "replace"))
    except Exception:
        return None

    def _call(v: Any) -> Any:
        return v() if callable(v) else v

    def kind_of(node: Any) -> str:
        t = getattr(node, "type", None)
        return t if isinstance(t, str) else _call(node.kind)

    def children_of(node: Any) -> List[Any]:
        ch = getattr(node, "children", None)
        if isinstance(ch, list):
            return ch
        return [node.child(i) for i in range(_call(node.child_count))]

    def text_of(node: Any) -> str:
        raw = getattr(node, "text", None)
        if isinstance(raw, bytes):
            return raw.decode("utf-8", "replace")
        return data[_call(node.start_byte):_call(node.end_byte)].decode("utf-8", "replace")

    def line_of(node: Any) -> int:
        pt = _call(node.start_point if hasattr(node, "start_point")
                   else node.start_position)
        return int(pt[0] if isinstance(pt, tuple) else _call(pt.row)) + 1

    text = data.decode("utf-8", "replace")
    units = [_file_unit(rel, "code", policy, text, lang)]
    units[0]["kind"] = "module"

    def name_of(node: Any) -> Optional[str]:
        n = node.child_by_field_name("name")
        if n is not None:
            return text_of(n)
        for c in children_of(node):
            if kind_of(c) in ("name", "identifier", "type_identifier", "field_identifier"):
                return text_of(c)
        return None

    def calls_in(node: Any) -> List[str]:
        out = []
        stack = [node]
        while stack:
            n = stack.pop()
            ch = children_of(n)
            if kind_of(n) in _TS_CALL:
                f = n.child_by_field_name("function")
                if f is None and ch:
                    f = ch[0]
                if f is not None:
                    t = text_of(f)
                    out.append(t.replace("(", " ").split()[0].split(".")[-1].split("::")[-1])
            stack.extend(ch)
        return sorted({c for c in out if c.isidentifier()})

    def walk(node: Any) -> None:
        for child in children_of(node):
            k = kind_of(child)
            if k in _TS_DEF:
                nm = name_of(child)
                if nm:
                    src = data[_call(child.start_byte):_call(child.end_byte)] \
                        .decode("utf-8", "replace")
                    units.append({
                        "id": f"code:{rel}::{nm}", "kind": _TS_DEF[k],
                        "source_path": rel,
                        "category": "code", "policy": policy, "qualname": nm,
                        "lineno": line_of(child), "lang": lang,
                        "content_sha": _sha(src), "calls": calls_in(child)})
            walk(child)

    if tree is None:                          # parser returned no tree -> skip
        return None
    walk(_call(tree.root_node))
    return units


_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
# Code-fence delimiter (CommonMark minimum): ``` or ~~~, 3+ chars, indented up
# to 3 spaces. A closing fence is the same character, at least as long.
_FENCE = re.compile(r"^ {0,3}(`{3,}|~{3,})")


def _slug(title: str) -> str:
    s = re.sub(r"[^\w\s-]", "", title.strip().lower())
    return re.sub(r"[\s-]+", "-", s)[:60] or "section"


def _markdown_units(path: Path, rel: str, policy: str) -> List[Dict[str, Any]]:
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines(keepends=True)
    sections: List[Tuple[Optional[str], int, int]] = []
    cur_title: Optional[str] = None
    cur_start = 0
    fence = None                 # opening fence marker while inside a code block
    for i, line in enumerate(lines):
        raw = line.rstrip("\n")
        fmatch = _FENCE.match(raw)
        if fmatch:
            marker = fmatch.group(1)
            if fence is None:
                fence = marker
            elif marker[0] == fence[0] and len(marker) >= len(fence):
                fence = None     # a foreign marker inside a block closes nothing
            continue
        if fence is not None:
            continue             # inside a fence: '# ...' is code, not a heading
        m = _HEADING.match(raw)
        if m:
            if cur_title is not None or i > 0:
                sections.append((cur_title, cur_start, i))
            cur_title, cur_start = m.group(2), i
    sections.append((cur_title, cur_start, len(lines)))

    units: List[Dict[str, Any]] = []
    seen: Dict[str, int] = {}
    for title, start, end in sections:
        chunk = "".join(lines[start:end]).strip()
        if not chunk:
            continue
        if title is None:
            qual = "_preamble"
        else:
            base = _slug(title)
            n = seen.get(base, 0)
            seen[base] = n + 1
            qual = base if n == 0 else f"{base}-{n}"
        units.append({"id": f"doc:{rel}::{qual}", "kind": "section", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": qual,
                      "lineno": start + 1, "lang": "markdown", "content_sha": _sha(chunk)})
    return units or [_file_unit(rel, "doc", policy, text, "markdown")]


def _text_units(path: Path, rel: str, policy: str) -> List[Dict[str, Any]]:
    text = path.read_text(encoding="utf-8", errors="replace")
    blocks, idx, n = [], 0, 0
    for raw in re.split(r"\n\s*\n", text):
        block = raw.strip()
        if not block:
            continue
        n += 1
        blocks.append({"id": f"doc:{rel}::b{n}", "kind": "block", "source_path": rel,
                       "category": "doc", "policy": policy, "qualname": f"b{n}",
                       "lineno": 1, "lang": "text", "content_sha": _sha(block)})
    return blocks or [_file_unit(rel, "doc", policy, text, "text")]


# A run of one punctuation char (>=3) used as an RST title adornment (under/overline).
_RST_ADORN = re.compile(r"^([=\-~^\"#*+.:'`_])\1{2,}\s*$")


def _rst_units(path: Path, rel: str, policy: str) -> List[Dict[str, Any]]:
    """reStructuredText sections by adornment: a title is a line whose NEXT line is a
    punctuation run at least as long as the title (RST's underline rule), optionally
    preceded by a matching overline. Splits into sections like markdown (the `headings`
    chunker only sees '#' headings, so .rst used to collapse to one prose section); text
    before the first heading is _preamble. Falls back to one file unit if no headings."""
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines(keepends=True)
    raw = [ln.rstrip("\n") for ln in lines]
    heads: List[Tuple[str, int]] = []           # (title, start line index)
    i = 0
    while i < len(raw):
        line = raw[i]
        if (_RST_ADORN.match(line) and i + 2 < len(raw) and raw[i + 1].strip()
                and _RST_ADORN.match(raw[i + 2]) and raw[i + 2][0] == line[0]):
            heads.append((raw[i + 1].strip(), i))      # overline + title + underline
            i += 3
            continue
        if (line.strip() and i + 1 < len(raw) and _RST_ADORN.match(raw[i + 1])
                and len(raw[i + 1].strip()) >= len(line.strip())):
            heads.append((line.strip(), i))            # title + underline
            i += 2
            continue
        i += 1
    if not heads:
        return [_file_unit(rel, "doc", policy, text, "rst")]

    units: List[Dict[str, Any]] = []
    seen: Dict[str, int] = {}
    if heads[0][1] > 0:
        pre = "".join(lines[:heads[0][1]]).strip()
        if pre:
            units.append({"id": f"doc:{rel}::_preamble", "kind": "section",
                          "source_path": rel, "category": "doc", "policy": policy,
                          "qualname": "_preamble", "lineno": 1, "lang": "rst",
                          "content_sha": _sha(pre)})
    bounds = [h[1] for h in heads] + [len(lines)]
    for idx, (title, start) in enumerate(heads):
        chunk = "".join(lines[start:bounds[idx + 1]]).strip()
        if not chunk:
            continue
        base = _slug(title)
        c = seen.get(base, 0)
        seen[base] = c + 1
        qual = base if c == 0 else f"{base}-{c}"
        units.append({"id": f"doc:{rel}::{qual}", "kind": "section", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": qual,
                      "lineno": start + 1, "lang": "rst", "content_sha": _sha(chunk)})
    return units or [_file_unit(rel, "doc", policy, text, "rst")]


# A line that begins with a recognizable timestamp -> the start of a log event.
_LOG_TS = re.compile(
    r"^\s*\[?"
    r"(?:\d{4}-\d{2}-\d{2}[ T]\d{2}:\d{2}:\d{2}"          # ISO: 2026-06-18 10:00:00 / T
    r"|\d{2}:\d{2}:\d{2}"                                  # bare: 10:00:00
    r"|[A-Z][a-z]{2}\s+\d{1,2}\s+\d{2}:\d{2}:\d{2})")      # syslog: Jun 18 10:00:00


def _log_units(path: Path, rel: str, policy: str, group_lines: int = 50) -> List[Dict[str, Any]]:
    """Group a .log into bounded episodes: timestamped lines confirm it is a log, then
    consecutive lines are bundled into windows of `group_lines` (one `block` unit each,
    `e{n}`). The window keeps a long log from becoming one huge node AND from exploding
    into one node per line; continuation lines (stack traces) ride along in their
    window. Carries `text` so the builder reads the slice once. A file with no
    timestamped line is not really a log -> paragraph blocks (_text_units)."""
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines(keepends=True)
    if not any(_LOG_TS.match(ln) for ln in lines):
        return _text_units(path, rel, policy)
    units: List[Dict[str, Any]] = []
    n = 0
    step = max(1, group_lines)
    for start in range(0, len(lines), step):
        chunk = "".join(lines[start:start + step]).strip()
        if not chunk:
            continue
        n += 1
        units.append({"id": f"doc:{rel}::e{n}", "kind": "block", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": f"e{n}",
                      "lineno": start + 1, "lang": "log", "content_sha": _sha(chunk),
                      "text": chunk})
    return units or [_file_unit(rel, "doc", policy, text, "log")]


def _json_qual(path: str) -> str:
    """A stable, <=48-char qualname for a nested JSON key path. A long path keeps a
    readable head plus a hash of the FULL path, so distinct deep paths never collide
    (unlike a plain truncation) and the id stays stable across runs regardless of entry
    order."""
    return (path if len(path) <= 48
            else path[:39] + "-" + hashlib.sha256(path.encode()).hexdigest()[:8])


def _has_child_container(v: object) -> bool:
    """True if a dict/list holds at least one nested dict/list. A large but FLAT
    container (e.g. a list of 10k numbers) has none, so it is kept as one record rather
    than exploded into a node per scalar — recursion targets nested STRUCTURE."""
    if isinstance(v, dict):
        return any(isinstance(x, (dict, list)) for x in v.values())
    if isinstance(v, list):
        return any(isinstance(x, (dict, list)) for x in v)
    return False


def _json_descend(rel: str, policy: str, parent: str, container: object, depth: int,
                  max_depth: int, recurse_min: int, units: List[Dict[str, Any]], cap: int) -> None:
    """Emit record units for the children of a large nested container, recursing into
    children that are themselves large nested structures (depth-limited). A child that
    is a leaf or a small/flat container becomes one record keyed by its full path
    (a.b.c / a.b[0]). Bounded by `cap` total units."""
    pairs = (container.items() if isinstance(container, dict)
             else enumerate(container) if isinstance(container, list) else [])
    for k, v in pairs:
        if len(units) >= cap:
            return
        path = f"{parent}.{k}" if isinstance(container, dict) else f"{parent}[{k}]"
        if (isinstance(v, (dict, list)) and depth < max_depth and _has_child_container(v)
                and len(json.dumps(v, ensure_ascii=False)) > recurse_min):
            _json_descend(rel, policy, path, v, depth + 1, max_depth, recurse_min, units, cap)
        else:
            qual = _json_qual(path)
            frag = json.dumps(v, ensure_ascii=False, sort_keys=True)
            units.append({"id": f"data:{rel}::{qual}", "kind": "record", "source_path": rel,
                          "category": "data", "policy": policy, "qualname": qual,
                          "lineno": 1, "lang": "json", "content_sha": _sha(frag)})


def _data_units(path: Path, rel: str, policy: str, max_depth: int = 4,
                recurse_min: int = 2048, cap: int = 500) -> List[Dict[str, Any]]:
    """JSON/YAML records. Each top-level entry is one record, EXCEPT a large nested
    container (serialized JSON over `recurse_min`, holding nested structure) is split
    into sub-records by key path so deep structure is not lost (recursive chunker,
    Stage 11). A small or flat value keeps the original one-record-per-entry shape and
    hash, so ordinary data files are unchanged. Total units per file are capped at
    `cap` (json_max_nodes); `max_depth` bounds recursion. A scalar root or parse error
    falls back to one file unit."""
    text = path.read_text(encoding="utf-8", errors="replace")
    try:
        obj = yaml.safe_load(text)            # parses JSON and YAML
    except Exception:
        return [_file_unit(rel, "data", policy, text)]
    if isinstance(obj, dict):
        items: List[Tuple[object, object]] = list(obj.items())
    elif isinstance(obj, list):
        items = list(enumerate(obj))
    else:
        return [_file_unit(rel, "data", policy, text)]
    units: List[Dict[str, Any]] = []
    for key, val in items:
        if len(units) >= cap:
            break
        if (isinstance(val, (dict, list)) and max_depth > 1 and _has_child_container(val)
                and len(json.dumps(val, ensure_ascii=False)) > recurse_min):
            _json_descend(rel, policy, str(key), val, 2, max_depth, recurse_min, units, cap)
        else:
            frag = json.dumps({str(key): val}, ensure_ascii=False, sort_keys=True)
            units.append({"id": f"data:{rel}::{str(key)[:48]}", "kind": "record",
                          "source_path": rel, "category": "data", "policy": policy,
                          "qualname": str(key)[:48], "lineno": 1, "lang": "json",
                          "content_sha": _sha(frag)})
    return units or [_file_unit(rel, "data", policy, text)]


def _ndjson_units(path: Path, rel: str, policy: str, cap: int = 500) -> List[Dict[str, Any]]:
    """Newline-delimited JSON (.ndjson): one `record` unit per line, since each line is
    an independent JSON value (yaml.safe_load, used by _data_units, cannot read this
    line-oriented form). A line's own id-ish field (id/key/name/_id) gives a stable
    qualname when present, else the 1-based line number (L{n}); lineno is the real
    source line. Unparseable lines are skipped; a file with none falls back to one file
    unit. Capped at the first `cap` records (json_max_nodes), like _data_units."""
    text = path.read_text(encoding="utf-8", errors="replace")
    units: List[Dict[str, Any]] = []
    seen: Dict[str, int] = {}
    n = 0
    for i, line in enumerate(text.splitlines(), 1):
        s = line.strip()
        if not s:
            continue
        try:
            obj = json.loads(s)
        except ValueError:
            continue
        n += 1
        key = None
        if isinstance(obj, dict):
            for k in ("id", "key", "name", "_id"):
                if obj.get(k) not in (None, ""):
                    key = str(obj[k])
                    break
        base = (key or f"L{n}")[:48]
        c = seen.get(base, 0)
        seen[base] = c + 1
        qual = base if c == 0 else f"{base}-{c}"
        frag = json.dumps(obj, ensure_ascii=False, sort_keys=True)
        units.append({"id": f"data:{rel}::{qual}", "kind": "record", "source_path": rel,
                      "category": "data", "policy": policy, "qualname": qual,
                      "lineno": i, "lang": "ndjson", "content_sha": _sha(frag)})
        if n >= cap:
            break
    return units or [_file_unit(rel, "data", policy, text, "ndjson")]


def _csv_units(path: Path, rel: str, policy: str) -> List[Dict[str, Any]]:
    """One STRUCTURAL unit per CSV/TSV file: a table is data, not prose (like XLSX), so
    the stored text describes it — column headers, row count, a few sample rows — not
    every cell. One unit per file keeps a 10k-row CSV from exploding into 10k nodes;
    deep per-row chunking is a recursive/data concern, not this. The delimiter is the
    tab for .tsv, else sniffed (comma fallback). Carries `text` so the builder
    summarizes without re-reading the file."""
    import csv
    text = path.read_text(encoding="utf-8", errors="replace")
    if path.suffix.lower() == ".tsv":
        delim = "\t"
    else:
        try:
            delim = csv.Sniffer().sniff(text[:2048], delimiters=",;\t|").delimiter
        except csv.Error:
            delim = ","
    rows = [r for r in csv.reader(text.splitlines(), delimiter=delim)
            if any(c.strip() for c in r)]
    if not rows:
        return [_file_unit(rel, "data", policy, text, "csv")]
    header = rows[0]
    cols = ", ".join(h.strip() for h in header if h.strip())
    sample = [" | ".join(c for c in r) for r in rows[1:4]]
    desc = (f"CSV table '{Path(rel).name}': {len(rows) - 1} rows x {len(header)} columns.\n"
            f"Columns: {cols}\nSample rows:\n" + "\n".join(sample))
    return [{"id": f"data:{rel}", "kind": "sheet", "source_path": rel,
             "category": "data", "policy": policy, "qualname": "", "lineno": 1,
             "lang": "csv", "content_sha": _sha(desc), "text": desc}]


# --- captured sessions (Stage 9) ---------------------------------------------
# The dump writer (lifecycle.py) and this chunker SHARE the role markers below — the
# format contract; lifecycle imports these helpers so the two never drift apart.
_SESSION_ROLE_RE = re.compile(r"^=== (Human|Assistant) ===\s*$")


def session_role_marker(role: str) -> str:
    """Turn delimiter (`=== Human ===` / `=== Assistant ===`) shared by the session
    dump writer and the session chunker."""
    return f"=== {role} ==="


def session_attachment_marker(n: int, label: str = "") -> str:
    """One marker per omitted attachment (tool call/result, image, file), numbered
    sequentially — so several attachments in one message stay distinct (e.g. a chat
    export where a user turn carries many files), not collapsed into a count. Writer-
    only in practice: the chunker keeps it inside the turn (it splits on role markers,
    not these)."""
    return f"== Attachment {n}: {label} ==" if label else f"== Attachment {n} =="


def _session_units(path: Path, rel: str, policy: str) -> List[Dict[str, Any]]:
    """One unit per conversation turn in a captured session dump. Turns are split on
    the role markers the dump writer emits; the leading frontmatter (everything before
    the first marker) is skipped. Each turn is a `section` so it is episodic — a long
    chat is ingested at conversational granularity and consolidation can summarize /
    compact piled-up sessions (THEORY §13). Falls back to one file unit if no markers
    are present (e.g. a hand-dropped chat export the writer did not format)."""
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines(keepends=True)
    starts = [i for i, ln in enumerate(lines) if _SESSION_ROLE_RE.match(ln.rstrip("\n"))]
    if not starts:
        return [_file_unit(rel, "doc", policy, text, "session")]
    bounds = starts + [len(lines)]
    units = []
    for n, (s, e) in enumerate(zip(bounds, bounds[1:]), 1):
        chunk = "".join(lines[s:e]).strip()
        if not chunk:
            continue
        units.append({"id": f"doc:{rel}::m{n}", "kind": "section", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": f"m{n}",
                      "lineno": s + 1, "lang": "session", "content_sha": _sha(chunk)})
    return units or [_file_unit(rel, "doc", policy, text, "session")]


# --- external chat exports (Stage 11) ----------------------------------------
# A structured chat log (JSON array / a {messages: [...]} object / NDJSON) of message
# objects -> one episodic `section` per message, with role/time/thread folded into the
# unit text (so the builder summarizes WITH attribution) and a weak `follows` edge to
# the previous turn IN THE SAME THREAD (conversation adjacency: roadmap 4.2). This is the
# common OpenAI/Anthropic `messages` shape and tolerant synonyms; our own flat dump
# (=== Human === markers) is handled by _session_units, reached via _has_role_markers.
_CHAT_ROLE_KEYS = ("role", "author", "speaker", "from", "sender", "name")
_CHAT_TEXT_KEYS = ("content", "text", "message", "body", "value")
_CHAT_TIME_KEYS = ("timestamp", "time", "ts", "created_at", "create_time", "date")
_CHAT_ID_KEYS = ("id", "message_id", "msg_id", "uuid")
_CHAT_THREAD_KEYS = ("conversation_id", "thread_id", "session_id", "channel", "chat_id",
                     "conversation")
_CHAT_CONTAINER_KEYS = ("messages", "conversation", "conversations", "turns", "history", "log")


def _msg_get(d: Dict[str, Any], keys: Tuple[str, ...]) -> Optional[object]:
    """First present, non-empty value among `keys` (synonym-tolerant field lookup)."""
    for k in keys:
        v = d.get(k)
        if v not in (None, ""):
            return v
    return None


def _chat_text(content: object) -> str:
    """Flatten a message's content to text. A string is itself; a list of blocks (the
    OpenAI/Anthropic shape) joins its text blocks and replaces each non-text block
    (tool_use, image) with a numbered attachment marker, the same convention the session
    dump uses; a {parts: [...]} / {text: ...} dict is unwrapped."""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts, att = [], 0
        for blk in content:
            if isinstance(blk, str):
                parts.append(blk)
            elif isinstance(blk, dict) and blk.get("text"):
                parts.append(str(blk["text"]))
            elif isinstance(blk, dict) and isinstance(blk.get("content"), str):
                parts.append(blk["content"])
            else:
                att += 1
                kind = blk.get("type") if isinstance(blk, dict) else "attachment"
                parts.append(session_attachment_marker(att, str(kind or "attachment")))
        return "\n".join(p for p in parts if p)
    if isinstance(content, dict):
        if isinstance(content.get("parts"), list):
            return _chat_text(content["parts"])
        if content.get("text"):
            return str(content["text"])
    return "" if content is None else str(content)


def _chat_messages(obj: object) -> Optional[List[Dict[str, Any]]]:
    """Extract the message list from a parsed chat export, or None if it does not look
    like one. Accepts a bare list or a dict carrying a messages-ish list; requires at
    least two dict records, most of which have a role-ish AND a text-ish field — so an
    ordinary JSON array of records is NOT misread as a chat."""
    seq: Optional[List[Any]] = None
    if isinstance(obj, list):
        seq = obj
    elif isinstance(obj, dict):
        for k in _CHAT_CONTAINER_KEYS:
            if isinstance(obj.get(k), list):
                seq = obj[k]
                break
    if not seq:
        return None
    msgs = [m for m in seq if isinstance(m, dict)]
    if len(msgs) < 2 or len(msgs) < 0.6 * len(seq):
        return None
    roled = sum(1 for m in msgs if _msg_get(m, _CHAT_ROLE_KEYS) is not None)
    texted = sum(1 for m in msgs if _msg_get(m, _CHAT_TEXT_KEYS) is not None)
    if roled < 0.6 * len(msgs) or texted < 0.6 * len(msgs):
        return None
    return msgs


def _chat_qual(mid: str) -> str:
    q = re.sub(r"[^\w.-]+", "_", mid).strip("_")[:48]
    return q or "m"


def _chat_units(path: Path, rel: str, policy: str) -> Optional[List[Dict[str, Any]]]:
    """One `section` unit per message of a structured chat export (JSON or NDJSON).
    Returns None when the file is not a recognizable chat -> the caller falls back to the
    json / ndjson record chunker. Each message carries `follows` (the previous turn in
    the same thread) so consecutive turns form a weak conversational chain in the graph;
    its id is the message's own id field when present (stable across re-exports), else
    m{seq}. role/timestamp/thread are folded into the unit text for the summary."""
    text = path.read_text(encoding="utf-8", errors="replace")
    try:
        obj: object = json.loads(text)
    except ValueError:                       # maybe NDJSON: one message object per line
        rows = []
        for line in text.splitlines():
            s = line.strip()
            if not s:
                continue
            try:
                rows.append(json.loads(s))
            except ValueError:
                return None                  # not clean NDJSON -> let ndjson chunker handle it
        obj = rows
    msgs = _chat_messages(obj)
    if not msgs:
        return None

    units: List[Dict[str, Any]] = []
    seen: Dict[str, int] = {}
    last_by_thread: Dict[str, str] = {}
    for n, m in enumerate(msgs, 1):
        role = str(_msg_get(m, _CHAT_ROLE_KEYS) or "unknown")
        body = _chat_text(_msg_get(m, _CHAT_TEXT_KEYS))
        ts = _msg_get(m, _CHAT_TIME_KEYS)
        thread = _msg_get(m, _CHAT_THREAD_KEYS)
        mid = _msg_get(m, _CHAT_ID_KEYS)
        base = _chat_qual(str(mid)) if mid is not None else f"m{n}"
        c = seen.get(base, 0)
        seen[base] = c + 1
        qual = base if c == 0 else f"{base}-{c}"
        header = (f"[{role}]" + (f" {ts}" if ts else "")
                  + (f" (thread {thread})" if thread is not None else ""))
        full = (header + "\n" + body).strip()
        unit: Dict[str, Any] = {"id": f"doc:{rel}::{qual}", "kind": "section", "source_path": rel,
                "category": "doc", "policy": policy, "qualname": qual, "lineno": 1,
                "lang": "chat", "content_sha": _sha(full), "text": full}
        tkey = str(thread) if thread is not None else "_"
        prev = last_by_thread.get(tkey)
        if prev:
            unit["follows"] = prev           # weak adjacency to the previous turn in-thread
        last_by_thread[tkey] = unit["id"]
        units.append(unit)
    return units or None


def _has_role_markers(path: Path) -> bool:
    """True if a prose file is actually a session/chat dump in the shared flat format
    (=== Human === / === Assistant ===). Lets a dump dropped into an ordinary
    mirror/absorb source route to the session chunker instead of being read as plain
    markdown/text. Cheap: only the head is scanned."""
    try:
        head = path.read_text(encoding="utf-8", errors="replace")[:4096]
    except OSError:
        return False
    return any(_SESSION_ROLE_RE.match(ln) for ln in head.splitlines())


# --- binary document formats (optional pure-Python libs; graceful skip) ------
# Each returns units that CARRY the extracted `text`, so the builder can summarize
# without re-opening the binary. Returns None when the library is absent or the
# file can't be read (e.g. a scanned PDF with no text layer) -> the file is skipped.

def _pdf_units(path: Path, rel: str, policy: str) -> Optional[List[Dict[str, Any]]]:
    try:
        from pypdf import PdfReader
    except Exception:
        return None
    try:
        reader = PdfReader(str(path))
    except Exception:
        return None
    units = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        if not text:
            continue
        units.append({"id": f"doc:{rel}::p{i}", "kind": "page", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": f"p{i}",
                      "lineno": 1, "lang": "pdf", "content_sha": _sha(text), "text": text})
    return units or None          # no extractable text (scanned?) -> skip


def _docx_units(path: Path, rel: str, policy: str) -> Optional[List[Dict[str, Any]]]:
    try:
        import docx                                  # python-docx
    except Exception:
        return None
    try:
        document = docx.Document(str(path))
    except Exception:
        return None
    sections: List[Tuple[Optional[str], List[str]]] = []
    title: Optional[str] = None
    body: List[str] = []                             # chunk by heading, like markdown
    for para in document.paragraphs:
        txt = para.text.strip()
        style = (para.style.name or "") if para.style else ""
        if style.lower().startswith("heading") or style.lower() == "title":
            if title is not None or body:
                sections.append((title, body))
            title, body = txt, []
        elif txt:
            body.append(txt)
    sections.append((title, body))

    units: List[Dict[str, Any]] = []
    seen: Dict[str, int] = {}
    for title, paras in sections:
        chunk = "\n".join(([title] if title else []) + paras).strip()
        if not chunk:
            continue
        base = _slug(title) if title else "_preamble"
        n = seen.get(base, 0)
        seen[base] = n + 1
        qual = base if n == 0 else f"{base}-{n}"
        units.append({"id": f"doc:{rel}::{qual}", "kind": "section", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": qual,
                      "lineno": 1, "lang": "docx", "content_sha": _sha(chunk), "text": chunk})
    return units or None


def _xlsx_units(path: Path, rel: str, policy: str) -> Optional[List[Dict[str, Any]]]:
    """One unit per SHEET. A spreadsheet is data, not prose, so the stored text is a
    STRUCTURAL description (sheet name, size, column headers, a few sample rows),
    not a dump of every cell."""
    try:
        import openpyxl
    except Exception:
        return None
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return None
    units = []
    for ws in wb.worksheets:
        rows = ws.iter_rows(values_only=True)
        header = next(rows, None) or ()
        cols = ", ".join(str(h) for h in header if h is not None)
        sample = []
        for r in rows:
            sample.append(" | ".join("" if v is None else str(v) for v in r))
            if len(sample) >= 3:
                break
        desc = (f"Spreadsheet sheet '{ws.title}': {ws.max_row or 0} rows x "
                f"{ws.max_column or 0} columns.\nColumns: {cols}\n"
                f"Sample rows:\n" + "\n".join(sample))
        units.append({"id": f"data:{rel}::{ws.title}", "kind": "sheet", "source_path": rel,
                      "category": "data", "policy": policy, "qualname": ws.title,
                      "lineno": 1, "lang": "xlsx", "content_sha": _sha(desc), "text": desc})
    wb.close()
    return units or None


def _pptx_units(path: Path, rel: str, policy: str) -> Optional[List[Dict[str, Any]]]:
    """One unit per SLIDE via python-pptx (optional pure-Python lib; graceful skip like
    PDF/DOCX/XLSX). Concatenates each slide's shape text and carries it as `text` so the
    builder summarizes without opening the binary. Returns None when the lib is absent
    or the file is unreadable -> the file is skipped, never an error. (Legacy .ppt has
    no reliable pure-Python reader and stays in BINARY_EXT.)"""
    try:
        from pptx import Presentation                    # python-pptx
    except Exception:
        return None
    try:
        prs = Presentation(str(path))
    except Exception:
        return None
    units = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [shape.text.strip() for shape in slide.shapes
                 if shape.has_text_frame and shape.text.strip()]
        body = "\n".join(parts).strip()
        if not body:
            continue
        units.append({"id": f"doc:{rel}::s{i}", "kind": "section", "source_path": rel,
                      "category": "doc", "policy": policy, "qualname": f"s{i}",
                      "lineno": 1, "lang": "pptx", "content_sha": _sha(body), "text": body})
    return units or None


CHUNKERS: Dict[str, Callable[..., Any]] = {"python": _python_units, "treesitter": _treesitter_units,
            "headings": _markdown_units, "paragraphs": _text_units, "json": _data_units,
            "ndjson": _ndjson_units, "csv": _csv_units, "rst": _rst_units, "log": _log_units,
            "pdf": _pdf_units, "docx": _docx_units, "xlsx": _xlsx_units, "pptx": _pptx_units,
            "session": _session_units}


# --------------------------------------------------------------------------- #
# Captured sessions as a source (Stage 9)
# --------------------------------------------------------------------------- #

def session_dir(project_root: Path, config: Dict[str, Any], amg_root: Optional[Path]) -> Optional[Path]:
    """Resolve the captured-sessions folder. `config['sessions']` is an optional
    project-relative override; by default it DERIVES as <store>/sessions from the
    resolved root, so it is correct under any agent dir (.claude / .agents / ...) with
    no installer dependency. Shared by the dump writer (lifecycle) and the chunker —
    one source of truth, so writer and reader can never diverge."""
    s = config.get("sessions")
    if s:
        return (project_root / str(s)).resolve()
    return (Path(amg_root).resolve() / "sessions") if amg_root else None


def _iter_session_files(base: Path) -> Iterable[Path]:
    """Yield files under the sessions dir, filtering only junk BELOW it. The dir lives
    inside the store (under the ignored agent dir), so — unlike iter_source_files — its
    prefix is NOT subject to DEFAULT_IGNORE_DIRS / .gitignore: it is an opted-in
    AMG-internal source. Without this, iter would silently drop every dump (audit 1.18)."""
    if not base.exists():
        return
    for p in base.rglob("*"):
        if p.is_file() and not any(part in DEFAULT_IGNORE_DIRS
                                   for part in p.relative_to(base).parts):
            yield p


def _extract_sessions(project_root: Path, config: Dict[str, Any],
                      amg_root: Optional[Path]) -> List[Dict[str, Any]]:
    """Captured session dumps as an ordinary source, routed to the session chunker.
    Policy comes from `session_policy` (absorb by default)."""
    base = session_dir(project_root, config, amg_root)
    if base is None:
        return []
    policy = str(config.get("session_policy", "absorb"))
    out: List[Dict[str, Any]] = []
    for p in _iter_session_files(base):
        try:
            rel = p.relative_to(project_root).as_posix()
        except ValueError:                # an override outside the project tree
            continue
        out += _session_units(p, rel, policy)
    return out


# --------------------------------------------------------------------------- #
# Driver
# --------------------------------------------------------------------------- #

def extract(project_root: Path, config: Dict[str, Any], amg_root: Optional[Path] = None) -> List[Dict[str, Any]]:
    gitignore = load_gitignore(project_root) if config.get("respect_gitignore", True) else []
    ignore_dirs = _effective_ignore_dirs(amg_root)
    overrides = load_overrides(amg_root)
    log_lines = int(config.get("log_group_lines", 50) or 50)   # episode window for .log
    j_depth = int(config.get("json_max_depth", 4) or 4)        # recursive-JSON tunables
    j_min = int(config.get("json_recurse_min_chars", 2048) or 2048)
    j_cap = int(config.get("json_max_nodes", 500) or 500)
    units: List[Dict[str, Any]] = []
    for base_rel, policy in resolve_sources(config):
        extra = _excludes_for_policy(config, policy)
        for path in iter_source_files(project_root, base_rel, extra, gitignore, ignore_dirs):
            rel = path.relative_to(project_root).as_posix()
            category, chunker, lang, _amb, _ov = _classify_path(path, rel, overrides)
            if chunker == "skip":
                continue
            if chunker == "treesitter":
                got = _treesitter_units(path, rel, policy, lang) if lang else None
                if got is None:                       # graceful degradation (no grammar)
                    got = [_file_unit(rel, "code", policy,
                                      path.read_text(encoding="utf-8", errors="replace"), lang)]
                units += got
            elif chunker in ("pdf", "docx", "xlsx", "pptx"):
                got = CHUNKERS[chunker](path, rel, policy)
                if got is None:                       # lib missing / unreadable -> skip
                    continue
                units += got
            elif chunker == "log":
                units += _log_units(path, rel, policy, group_lines=log_lines)
            elif chunker in ("json", "ndjson"):
                got = _chat_units(path, rel, policy)       # structured chat export?
                if got is None:
                    got = (_ndjson_units(path, rel, policy, cap=j_cap) if chunker == "ndjson"
                           else _data_units(path, rel, policy, max_depth=j_depth,
                                            recurse_min=j_min, cap=j_cap))
                units += got
            elif chunker in ("headings", "paragraphs") and _has_role_markers(path):
                units += _session_units(path, rel, policy)  # a dropped flat role-marker dump
            else:
                units += CHUNKERS[chunker](path, rel, policy)
    units += _extract_sessions(project_root, config, amg_root)   # opted-in store source
    return units


def _stats(project_root: Path, config: Dict[str, Any], amg_root: Optional[Path] = None) -> Dict[str, Any]:
    gitignore = load_gitignore(project_root) if config.get("respect_gitignore", True) else []
    ignore_dirs = _effective_ignore_dirs(amg_root)
    overrides = load_overrides(amg_root)
    from collections import Counter
    cat: Counter[str] = Counter()
    langs: Counter[str] = Counter()
    skipped = 0
    ambiguous: List[str] = []
    resolved: List[str] = []
    by_source: Dict[str, Dict[str, Any]] = {}
    seen_src: Dict[str, str] = {}               # rel -> first source policy (overlap detect, 1.29)
    overlaps: List[str] = []
    for base_rel, policy in resolve_sources(config):
        extra = _excludes_for_policy(config, policy)
        found = (project_root / base_rel).exists()
        matched = 0
        for path in iter_source_files(project_root, base_rel, extra, gitignore, ignore_dirs):
            matched += 1                        # passed every ignore filter
            rel = path.relative_to(project_root).as_posix()
            if seen_src.get(rel, policy) != policy:
                overlaps.append(rel)            # same file under two different-policy roots
            seen_src.setdefault(rel, policy)
            category, chunker, lang, amb, overridden = _classify_path(path, rel, overrides)
            if chunker == "skip":
                skipped += 1
                continue
            cat[category] += 1
            langs[lang or chunker] += 1
            if overridden:
                resolved.append(rel)            # ambiguity settled by amg-classifier
            elif amb:
                ambiguous.append(rel)           # still unresolved -> prose default
        by_source[base_rel] = {"policy": policy, "found": found, "files": matched}
    sessions = 0
    sess_base = session_dir(project_root, config, amg_root)
    if sess_base:
        for _p in _iter_session_files(sess_base):
            cat["doc"] += 1
            langs["session"] += 1
            sessions += 1
    try:
        from tree_sitter_language_pack import get_parser
        get_parser("json")                    # attempt a real grammar load
        ts = "available (function-level units + calls for all languages)"
    except Exception:
        ts = "unavailable (non-Python code -> one unit per file; Python unaffected)"

    import importlib
    pip_name = {"docx": "python-docx", "pptx": "python-pptx"}
    extraction = {}
    for fmt, mod in (("pdf", "pypdf"), ("docx", "docx"), ("xlsx", "openpyxl"), ("pptx", "pptx")):
        try:
            importlib.import_module(mod)
            extraction[fmt] = f"available ({mod})"
        except Exception:
            extraction[fmt] = f"unavailable (pip install {pip_name.get(mod, mod)}) -> files skipped"

    out = {"by_category": dict(cat), "by_language": dict(langs),
           "by_source": by_source,             # per-source file counts (0 / not-found visible)
           "skipped_binary": skipped, "sessions": sessions,
           "ambiguous_files": ambiguous[:20],
           "resolved_by_override": resolved[:20],
           "tree_sitter": ts, "text_extraction": extraction}
    if overlaps:
        out["overlapping_sources"] = sorted(set(overlaps))[:20]   # 1.29: not resolved silently
        out["overlap_hint"] = (
            f"{len(set(overlaps))} file(s) fall under both mirror_path and absorb_path; "
            "the last-listed policy wins (absorb/absorb_once over mirror). Narrow the "
            "roots or add an exclude to resolve.")
    if ambiguous:
        ov_path = ((amg_root / "work" / "classification-overrides.json").as_posix()
                   if amg_root else "work/classification-overrides.json")
        out["classifier_hint"] = (
            f"{len(ambiguous)} ambiguous file(s) defaulted to prose. Spawn the "
            f"amg-classifier subagent on ambiguous_files, write its "
            f"{{path: {{category, language}}}} mapping to {ov_path}, then re-run.")
    return out


def main(argv: List[str]) -> int:
    project_root = Path(argv[1]).resolve() if len(argv) > 1 else Path.cwd()
    rest = [a for a in argv[2:] if not a.startswith("--")]
    if rest:
        amg_root = Path(rest[0]).resolve()
    else:
        # Resolve the store like every other CLI (reconcile/consolidate/notes) so a
        # global engine or a non-.claude agent dir still finds the project's graph.
        from graph_store import resolve_amg_root
        amg_root = resolve_amg_root(start=project_root)
    config = load_config(amg_root)
    if "--stats" in argv:
        print(json.dumps(_stats(project_root, config, amg_root), ensure_ascii=False, indent=2))
        return 0
    json.dump(extract(project_root, config, amg_root), sys.stdout, ensure_ascii=False, indent=2)
    print()
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
